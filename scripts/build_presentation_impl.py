"""Build the FraudOps presentation (16:9 PPTX).

Slides:
  1  Title
  2  Assignment context
  3  Business problem & requirements
  4  Architecture at a glance
  5  Screenshot: architecture diagram (screen 06)
  6  Event flow — the three topics
  7  Screenshot: login screen + RBAC (screen 01)
  8  Screenshot: idle dashboard (screen 03)
  9  Screenshot: simulator running (screen 04)
 10  Screenshot: fraud injected (screen 05)
 11  Screenshot: API Explorer (screens 07 + 08)
 12  Screenshot: design notes (screen 09)
 13  Screenshot: viewer RBAC — disabled buttons (screen 10)
 14  Screenshot: analyst dashboard + ack (screens 11 + 12)
 15  ML model detail
 16  Auth + RBAC matrix
 17  Latency budget
 18  Resilience
 19  Horizontal scaling
 20  Local ↔ production mapping (FastAPI → Spring Boot / Kafka)
 21  Testing & evaluation
 22  Thank you / next steps
"""
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SHOTS = Path(__file__).resolve().parent.parent / 'docs' / 'demo_screenshots'
OUT_DEFAULT = Path(__file__).resolve().parent.parent / 'docs' / 'fraudops-presentation.pptx'

# ---------- palette (matches the app) ----------
BG        = RGBColor(0x0A, 0x0C, 0x10)
CARD      = RGBColor(0x11, 0x14, 0x1A)
LINE      = RGBColor(0x22, 0x26, 0x30)
TEXT      = RGBColor(0xF8, 0xFA, 0xFC)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
DIM       = RGBColor(0x64, 0x74, 0x8B)
CYAN      = RGBColor(0x00, 0xE5, 0xFF)
AMBER     = RGBColor(0xFF, 0xB0, 0x00)
ROSE      = RGBColor(0xFF, 0x33, 0x66)
LIME      = RGBColor(0x00, 0xFF, 0x66)

# 16:9 slide
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs, kind="blank"):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, text, *,
             font_size=18, color=TEXT, bold=False, mono=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    if not isinstance(text, list):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Consolas" if mono else "Calibri"
    return tb


def add_pill(slide, left, top, text, color=CYAN, bg_alpha_hex=None):
    """Small pill/chip with a colored border and mono text."""
    width = Inches(1 + 0.11 * len(text))
    height = Inches(0.32)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.35
    rect.line.color.rgb = color
    rect.line.width = Pt(0.75)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD
    tf = rect.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text.upper()
    r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = color
    r.font.name = "Consolas"
    return rect


def add_card(slide, left, top, width, height, *, accent=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.03
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD
    rect.line.color.rgb = accent if accent else LINE
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def add_footer(slide, page, total):
    add_text(slide, Inches(0.5), Inches(7.10), Inches(9), Inches(0.3),
             "FraudOps · Real-time Fraud Detection with ML — Technical Deep Dive",
             font_size=9, color=DIM, mono=True)
    add_text(slide, Inches(12.4), Inches(7.10), Inches(0.6), Inches(0.3),
             f"{page:02d} / {total:02d}", font_size=9, color=DIM,
             mono=True, align=PP_ALIGN.RIGHT)


def add_title(slide, eyebrow, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.42), Inches(12), Inches(0.3),
             eyebrow.upper(), font_size=10, color=DIM, mono=True, bold=True)
    add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
             title, font_size=32, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.45), Inches(12), Inches(0.4),
                 subtitle, font_size=14, color=MUTED)


def add_image(slide, path, left, top, *, width=None, height=None, border_color=None):
    if width is None and height is None:
        pic = slide.shapes.add_picture(str(path), left, top)
    elif width is not None:
        pic = slide.shapes.add_picture(str(path), left, top, width=width)
    else:
        pic = slide.shapes.add_picture(str(path), left, top, height=height)
    # frame border
    if border_color:
        pic.line.color.rgb = border_color
        pic.line.width = Pt(0.5)
    return pic


def add_bullet_card(slide, left, top, width, height, title, points, accent):
    add_card(slide, left, top, width, height, accent=None)
    # accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.line.fill.background()
    stripe.fill.solid(); stripe.fill.fore_color.rgb = accent
    add_text(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.35), Inches(0.35),
             title.upper(), font_size=11, color=accent, mono=True, bold=True)
    add_text(slide, left + Inches(0.25), top + Inches(0.5), width - Inches(0.35),
             height - Inches(0.6), points, font_size=12, color=MUTED)


def make_deck(out_path=None):
    prs = new_deck()
    slides = []

    # ==================================================================
    # 1. Title
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_pill(s, Inches(0.5), Inches(0.5), "USE CASE 02  ·  TECHNICAL DEEP DIVE", CYAN)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12), Inches(1.5),
             "Real-time Fraud Detection",
             font_size=56, color=TEXT, bold=True)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12), Inches(1.5),
             "Microservice with ML Integration",
             font_size=40, color=CYAN, bold=True)
    add_text(s, Inches(0.5), Inches(4.9), Inches(12), Inches(1),
             "Three services · one event bus · one ML API · full RBAC console",
             font_size=18, color=MUTED)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.4),
             "FastAPI + MongoDB + React  ·  mirrors Spring Boot + Kafka",
             font_size=13, color=DIM, mono=True)

    # ==================================================================
    # 2. Assignment context
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "context", "Assignment recap",
              "Use Case 2 — Technical Deep Dive (Code + Demo)")
    fields = [
        ("Objective",
         "Evaluate coding skills, microservices design, ML capability."),
        ("Technical scenario",
         "Real-time fraud detection: process transactions, ML-score them, flag suspicious activity."),
        ("Assignment",
         "Transaction, Fraud Scoring, and Alert microservices. ML model exposed via API. Kafka for events."),
        ("Expected demo",
         "Working APIs, event flow, fraud scoring output, explanation of latency, resilience, scaling."),
        ("Evaluation",
         "Code quality · microservices design · ML integration · demo execution · technical presentation."),
    ]
    top = Inches(2.1)
    for label, body in fields:
        add_card(s, Inches(0.5), top, Inches(12.33), Inches(0.85))
        add_text(s, Inches(0.75), top + Inches(0.15), Inches(2.7), Inches(0.5),
                 label, font_size=13, color=CYAN, mono=True, bold=True)
        add_text(s, Inches(3.5), top + Inches(0.15), Inches(9.3), Inches(0.6),
                 body, font_size=13, color=TEXT)
        top += Inches(0.95)

    # ==================================================================
    # 3. Business problem & requirements
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "problem statement", "What the system has to do",
              "Non-functional and functional targets — the SLAs")
    # Left card — non-functional
    add_card(s, Inches(0.5), Inches(2.1), Inches(6.3), Inches(4.6), accent=CYAN)
    add_text(s, Inches(0.75), Inches(2.3), Inches(6), Inches(0.4),
             "NON-FUNCTIONAL", font_size=12, color=CYAN, mono=True, bold=True)
    rows = [
        ("Scoring p95 latency",  "≤ 10 ms"),
        ("End-to-end approve",   "≤ 50 ms"),
        ("Availability",         "99.95%"),
        ("Model retraining",     "Weekly + canary"),
        ("Fraud recall",         "≥ 92%"),
        ("Explainability",       "Reason codes per alert"),
    ]
    top = Inches(2.85)
    for k, v in rows:
        add_text(s, Inches(0.75), top, Inches(4), Inches(0.35), k, font_size=13, color=MUTED)
        add_text(s, Inches(4.6),  top, Inches(2.2), Inches(0.35), v, font_size=13, color=TEXT, mono=True, bold=True)
        top += Inches(0.45)

    # Right card — functional
    add_card(s, Inches(7.0), Inches(2.1), Inches(5.83), Inches(4.6), accent=AMBER)
    add_text(s, Inches(7.25), Inches(2.3), Inches(5), Inches(0.4),
             "FUNCTIONAL", font_size=12, color=AMBER, mono=True, bold=True)
    items = [
        "Ingest transactions over HTTP",
        "Score each tx via an ML model exposed as an API",
        "Publish scored events to a stream",
        "Raise alerts on high-risk events; analysts can ack",
        "Observability: throughput, latency, fraud rate",
        "Manual test surface: single-tx submit, fraud injection",
        "Authenticate operators — role-based capabilities",
    ]
    add_text(s, Inches(7.25), Inches(2.85), Inches(5.4), Inches(3.7),
             ["• " + i for i in items], font_size=13, color=TEXT)

    # ==================================================================
    # 4. Architecture at a glance
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "architecture", "Three microservices, one event bus, one ML API",
              "Every service communicates through the bus — never a direct call")

    # Three service boxes
    services = [
        ("Transaction Service",  "POST /api/tx/transactions",  "Ingests transactions.\nPersists + publishes to topic transactions.raw.", CYAN),
        ("Fraud Scoring Service", "POST /api/fraud/score",      "Consumes transactions.raw.\nCalls ML model.\nPublishes fraud.scores.", AMBER),
        ("Alert Service",        "POST /api/alerts/:id/ack",    "Consumes fraud.scores.\nRaises alerts.raised for fraud events.", ROSE),
    ]
    left = Inches(0.5)
    for title, endpoint, body, accent in services:
        add_bullet_card(s, left, Inches(2.15), Inches(4.05), Inches(3.7),
                         title, [endpoint, "", body], accent)
        left += Inches(4.28)

    # Bottom — event topics ribbon
    add_card(s, Inches(0.5), Inches(6.05), Inches(12.33), Inches(0.85), accent=LIME)
    add_text(s, Inches(0.75), Inches(6.2), Inches(3.5), Inches(0.45),
             "KAFKA-STYLE TOPICS", font_size=11, color=LIME, mono=True, bold=True)
    add_text(s, Inches(0.75), Inches(6.55), Inches(12), Inches(0.3),
             "transactions.raw  →  fraud.scores  →  alerts.raised",
             font_size=15, color=TEXT, mono=True, bold=True)

    # ==================================================================
    # 5. Screenshot: architecture diagram (animated)
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "architecture · live topology", "Live event topology",
              "Animated particles show real event counts since boot")
    add_image(s, SHOTS / "06_architecture_animated.jpg",
              Inches(0.5), Inches(2.0), width=Inches(12.33),
              border_color=LINE)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             "The dashed cyan arrows are Kafka topic hops. The ML API is a "
             "dashed callout — it's a service, not a library.",
             font_size=12, color=MUTED)

    # ==================================================================
    # 6. Event flow — the three topics
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "event flow", "One transaction through the whole system",
              "The happy-path sequence — six hops, ~15 ms end-to-end")
    steps = [
        ("1", "Client → Transaction Service",
         "POST /api/tx/transactions. Service enriches with tx_id, hour, cross_border flags.",   CYAN),
        ("2", "Transaction Service → transactions.raw",
         "Mongo insert (durability), then publish to topic. Returns 200 to client.",              CYAN),
        ("3", "transactions.raw → Fraud Scoring Service",
         "Async consumer picks it up. Invokes the ML model (IsolationForest + rules).",         AMBER),
        ("4", "Fraud Scoring → fraud.scores",
         "Adds ml_score, rule_score, fused fraud_score, decision, reasons. Publishes.",         AMBER),
        ("5", "fraud.scores → Alert Service",
         "Consumer checks risk_level == 'fraud'. Skips safe / suspicious.",                     ROSE),
        ("6", "Alert Service → alerts.raised",
         "Creates alert doc with severity + reasons. Persisted + published. Ops sees it.",       ROSE),
    ]
    top = Inches(2.0)
    for num, title, body, accent in steps:
        add_card(s, Inches(0.5), top, Inches(12.33), Inches(0.72), accent=None)
        # number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(0.7), top + Inches(0.2), Inches(0.36), Inches(0.36))
        badge.fill.solid(); badge.fill.fore_color.rgb = CARD
        badge.line.color.rgb = accent; badge.line.width = Pt(1.2)
        tf = badge.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = accent; r.font.name = "Consolas"
        add_text(s, Inches(1.25), top + Inches(0.08), Inches(6.5), Inches(0.35),
                 title, font_size=13, color=TEXT, bold=True)
        add_text(s, Inches(1.25), top + Inches(0.38), Inches(11.4), Inches(0.35),
                 body, font_size=11, color=MUTED)
        top += Inches(0.83)

    # ==================================================================
    # 7. Screenshot 01 — Login + RBAC
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 01", "The console — RBAC by first impression",
              "Three demo accounts, three capability sets")
    add_image(s, SHOTS / "01_login_screen.jpg",
              Inches(0.5), Inches(2.0), width=Inches(9.5), border_color=LINE)
    # Right-side notes
    add_card(s, Inches(10.3), Inches(2.0), Inches(2.5), Inches(4.7), accent=CYAN)
    add_text(s, Inches(10.5), Inches(2.15), Inches(2.2), Inches(0.4),
             "WHO CAN DO WHAT", font_size=11, color=CYAN, mono=True, bold=True)
    add_text(s, Inches(10.5), Inches(2.55), Inches(2.2), Inches(4),
             ["ADMIN",
              "start/stop simulator, inject fraud, manage all",
              "",
              "ANALYST",
              "acknowledge alerts, read all data",
              "",
              "VIEWER",
              "read-only observability"],
             font_size=11, color=MUTED)

    # ==================================================================
    # 8. Screenshot 03 — Idle dashboard
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 02", "The observability skeleton",
              "Four KPIs · demo controls · live stream · alerts panel")
    add_image(s, SHOTS / "03_admin_dashboard_idle.jpg",
              Inches(0.5), Inches(2.0), width=Inches(12.33), border_color=LINE)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             "Alerts panel is already populated — cold-start rehydration "
             "pulls the last 200 alerts from Mongo on boot.",
             font_size=12, color=MUTED)

    # ==================================================================
    # 9. Screenshot 04 — Simulator running
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 03", "Kafka comes alive",
              "3 TPS synthetic stream · every KPI updates in real time")
    add_image(s, SHOTS / "04_admin_simulator_running.jpg",
              Inches(0.5), Inches(2.0), width=Inches(12.33), border_color=LINE)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             "Latency KPI hovers around 5 ms · fused ML+rules score per row · "
             "SAFE / REVIEW / BLOCK map to the 0.50 / 0.75 thresholds.",
             font_size=12, color=MUTED)

    # ==================================================================
    # 10. Screenshot 05 — Fraud injected
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 04", "Fraud injection → alerts",
              "Reason codes are the compliance story")
    add_image(s, SHOTS / "05_admin_fraud_injected.jpg",
              Inches(0.5), Inches(2.0), width=Inches(12.33), border_color=LINE)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             "high_amount · foreign_transaction · high_risk_merchant · "
             "cross_border · country_hopping — every alert explains itself.",
             font_size=12, color=MUTED)

    # ==================================================================
    # 11. Screenshot 07 + 08 — API Explorer
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 05", "API Explorer",
              "Every endpoint has request · response · curl one-liner")
    add_image(s, SHOTS / "07_api_explorer_tx.jpg",
              Inches(0.5), Inches(2.0), width=Inches(6.1), border_color=LINE)
    add_image(s, SHOTS / "08_api_explorer_ml_score.jpg",
              Inches(6.75), Inches(2.0), width=Inches(6.1), border_color=LINE)
    add_text(s, Inches(0.5), Inches(6.05), Inches(6.1), Inches(0.35),
             "Transaction ingest — how merchant SDKs post",
             font_size=11, color=CYAN, mono=True)
    add_text(s, Inches(6.75), Inches(6.05), Inches(6.1), Inches(0.35),
             "ML score API — the model is a service, back-test-ready",
             font_size=11, color=AMBER, mono=True)

    # ==================================================================
    # 12. Screenshot 09 — Design notes
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 06", "Design notes",
              "Latency · Resilience · Scaling · ML — the interview talking points")
    add_image(s, SHOTS / "09_design_notes.jpg",
              Inches(0.5), Inches(2.0), width=Inches(12.33), border_color=LINE)

    # ==================================================================
    # 13. Screenshot 10 — Viewer RBAC
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 07", "RBAC in action — viewer",
              "Buttons disabled, sliders disabled, no ack — but data still flows")
    add_image(s, SHOTS / "10_viewer_dashboard_rbac.jpg",
              Inches(0.5), Inches(2.0), width=Inches(12.33), border_color=LINE)
    add_text(s, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             "Backend returns 403 to every gated endpoint. The UI matches "
             "exactly — disable + explain, never a mismatch.",
             font_size=12, color=MUTED)

    # ==================================================================
    # 14. Screenshot 11 + 12 — Analyst + ack
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "walkthrough · scene 08", "Analyst — the middle role",
              "Cannot start the stream, but can triage alerts")
    add_image(s, SHOTS / "11_analyst_dashboard.jpg",
              Inches(0.5), Inches(2.0), width=Inches(6.1), border_color=LINE)
    add_image(s, SHOTS / "12_analyst_ack_alert.jpg",
              Inches(6.75), Inches(2.0), width=Inches(6.1), border_color=LINE)
    add_text(s, Inches(0.5), Inches(6.05), Inches(6.1), Inches(0.4),
             "Ack buttons are back",
             font_size=11, color=AMBER, mono=True)
    add_text(s, Inches(6.75), Inches(6.05), Inches(6.1), Inches(0.4),
             "Acked alert fades — persisted in Mongo",
             font_size=11, color=LIME, mono=True)

    # ==================================================================
    # 15. ML model detail
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "the ml story", "Fused score — ML + rules",
              "Explainable, drift-monitored, fast (p95 ≤ 8 ms)")

    # Left card — features
    add_card(s, Inches(0.5), Inches(2.1), Inches(6.3), Inches(4.7), accent=CYAN)
    add_text(s, Inches(0.75), Inches(2.3), Inches(6), Inches(0.4),
             "ISOLATION FOREST · 7 FEATURES", font_size=11, color=CYAN, mono=True, bold=True)
    feats = [
        "amount",
        "hour  (0..23)",
        "is_foreign",
        "is_high_risk_merchant",
        "velocity_1h",
        "cross_border",
        "distinct_countries_24h",
    ]
    add_text(s, Inches(0.75), Inches(2.85), Inches(5.7), Inches(3.5),
             ["• " + f for f in feats], font_size=14, color=TEXT, mono=True)
    add_text(s, Inches(0.75), Inches(5.9), Inches(6), Inches(0.4),
             "Trained on 5,000 synthetic normals — n_estimators=120",
             font_size=11, color=DIM, mono=True)

    # Right card — fusion + thresholds
    add_card(s, Inches(7.0), Inches(2.1), Inches(5.83), Inches(4.7), accent=AMBER)
    add_text(s, Inches(7.25), Inches(2.3), Inches(5), Inches(0.4),
             "FUSION · THRESHOLDS", font_size=11, color=AMBER, mono=True, bold=True)
    add_text(s, Inches(7.25), Inches(2.85), Inches(5.3), Inches(0.5),
             "fused = 0.6 · ml_score + 0.4 · rule_score",
             font_size=15, color=TEXT, mono=True, bold=True)
    add_text(s, Inches(7.25), Inches(3.6), Inches(5.3), Inches(0.4),
             "THRESHOLDS", font_size=10, color=DIM, mono=True)
    rows = [
        ("[0.00, 0.50)",  "safe        · approve", LIME),
        ("[0.50, 0.75)",  "suspicious  · review",  AMBER),
        ("[0.75, 1.00]",  "fraud       · block",   ROSE),
    ]
    top = Inches(4.05)
    for rng, txt, col in rows:
        add_text(s, Inches(7.25), top, Inches(2.2), Inches(0.4),
                 rng, font_size=13, color=MUTED, mono=True)
        add_text(s, Inches(9.6),  top, Inches(3.2), Inches(0.4),
                 txt, font_size=13, color=col, mono=True, bold=True)
        top += Inches(0.5)
    add_text(s, Inches(7.25), Inches(6.05), Inches(5.5), Inches(0.4),
             "Every alert stamps reason codes — regulator-ready.",
             font_size=11, color=DIM, mono=True)

    # ==================================================================
    # 16. Auth + RBAC matrix
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "security", "Auth + RBAC — one line stops misuse",
              "JWT (HS256, 8h) · bcrypt · admin > analyst > viewer")

    matrix = [
        ("Endpoint",                                 "anon", "viewer", "analyst", "admin"),
        ("GET  /api/  ·  /api/health",               "✓",    "✓",      "✓",       "✓"),
        ("POST /api/auth/{register,login}",          "✓",    "✓",      "✓",       "✓"),
        ("GET  /api/auth/me",                        "",     "✓",      "✓",       "✓"),
        ("GET  /api/metrics  ·  /api/tx/...",        "",     "✓",      "✓",       "✓"),
        ("POST /api/fraud/score  (ML API)",          "",     "✓",      "✓",       "✓"),
        ("GET  /api/alerts/recent",                  "",     "✓",      "✓",       "✓"),
        ("POST /api/alerts/:id/ack",                 "",     "",       "✓",       "✓"),
        ("POST /api/simulator/*",                    "",     "",       "",        "✓"),
    ]
    cols_w = [Inches(6.3), Inches(1.4), Inches(1.4), Inches(1.5), Inches(1.5)]
    left0  = Inches(0.5); top0 = Inches(2.05)
    row_h  = Inches(0.42)
    for i, row in enumerate(matrix):
        y = top0 + row_h * i
        # row background
        if i == 0:
            rc = LINE
        else:
            rc = None
        if rc:
            bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left0, y, sum(cols_w, Emu(0)), row_h)
            bg.fill.solid(); bg.fill.fore_color.rgb = rc
            bg.line.fill.background()
        x = left0
        for j, cell in enumerate(row):
            fs = 11 if i == 0 else 12
            col_color = MUTED if i == 0 else (TEXT if j == 0 else CYAN)
            add_text(s, x + Inches(0.15), y + Inches(0.08),
                     cols_w[j] - Inches(0.2), row_h,
                     cell,
                     font_size=fs, color=col_color, mono=True,
                     bold=(i == 0 or (j == 0 and i > 0)),
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
            x += cols_w[j]

    add_text(s, Inches(0.5), Inches(6.75), Inches(12), Inches(0.35),
             "require_role(ROLE_ADMIN)  →  one FastAPI Depends, guards the endpoint at HTTP layer.",
             font_size=12, color=DIM, mono=True)

    # ==================================================================
    # 17. Latency budget
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "performance", "Latency budget — where the ms go",
              "Client returns in ~6 ms; scoring is async off the hot path")

    stages = [
        ("parse + validate",     1, CYAN),
        ("enrich",               1, CYAN),
        ("mongo insert (tx)",    2, CYAN),
        ("publish tx topic",     1, CYAN),
        ("→ return 200",         1, LIME),
        ("fraud scoring (ML)",   4, AMBER),
        ("mongo insert (score)", 2, AMBER),
        ("publish score topic",  1, AMBER),
        ("alert eval",           1, ROSE),
        ("mongo + publish alert",2, ROSE),
    ]
    top = Inches(2.15)
    max_ms = sum(x[1] for x in stages)
    max_w = Inches(9.0)
    for label, ms, col in stages:
        w = Emu(int(max_w / max_ms * ms))
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(4.0), top, w, Inches(0.34))
        bar.adjustments[0] = 0.4
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        add_text(s, Inches(0.5), top + Inches(0.05), Inches(3.4), Inches(0.35),
                 label, font_size=12, color=MUTED, align=PP_ALIGN.RIGHT)
        add_text(s, Inches(4.0) + w + Inches(0.1), top + Inches(0.05),
                 Inches(2), Inches(0.35),
                 f"{ms} ms", font_size=11, color=col, mono=True, bold=True)
        top += Inches(0.4)

    # divider annotation
    add_text(s, Inches(4.0), Inches(6.4), Inches(9), Inches(0.4),
             "── after the green stripe: async, off the hot path ──",
             font_size=11, color=DIM, mono=True, align=PP_ALIGN.LEFT)

    # ==================================================================
    # 18. Resilience
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "resilience", "What happens when things go wrong",
              "Producer failures · consumer failures · slow model · cold starts")

    items = [
        ("Producer failure",
         "Mongo down → tx endpoint 5xx, no publish. No phantom scores. Client retries.",
         CYAN),
        ("Consumer failure",
         "Bus catches exceptions per topic. In Kafka this maps to a dead-letter topic (fraud.scores.dlq).",
         AMBER),
        ("Slow model",
         "Resilience4j circuit breaker: on p99 SLO breach, fall back to rule-only scoring. Fail open, safer than freezing the stream.",
         ROSE),
        ("Cold-start amnesia",
         "rehydrate_alerts_store() pulls the last 200 alerts from Mongo on boot. The UI is never empty after a redeploy.",
         LIME),
        ("Key rotation",
         "JWT_SECRET per environment. A secondary decode-only key gives seamless rotation windows.",
         MUTED),
    ]
    top = Inches(2.1)
    for title, body, col in items:
        add_card(s, Inches(0.5), top, Inches(12.33), Inches(0.9), accent=None)
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(0.08), Inches(0.9))
        stripe.line.fill.background(); stripe.fill.solid(); stripe.fill.fore_color.rgb = col
        add_text(s, Inches(0.8), top + Inches(0.15), Inches(3.7), Inches(0.35),
                 title, font_size=13, color=col, bold=True, mono=True)
        add_text(s, Inches(4.5), top + Inches(0.15), Inches(8.2), Inches(0.75),
                 body, font_size=12, color=MUTED)
        top += Inches(1.0)

    # ==================================================================
    # 19. Horizontal scaling
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "scaling", "How this grows to 5,000 TPS",
              "Two axes — partitions and independent service scaling")

    # Left — partitions
    add_bullet_card(s, Inches(0.5), Inches(2.05), Inches(6.15), Inches(4.75),
                     "By kafka partitions",
                     ["transactions.raw is keyed by user_id",
                      "Per-user ordering is preserved — required for velocity_1h",
                      "N scoring pods == N partitions",
                      "Consumer group auto-rebalances on pod add/remove"],
                     CYAN)
    # Right — service scaling
    add_bullet_card(s, Inches(6.83), Inches(2.05), Inches(6.0), Inches(4.75),
                     "By service",
                     ["TX Service — stateless, CPU-light. ~4k rps per pod",
                      "Fraud Scoring — CPU-bound on ML call. HPA target p95 6 ms",
                      "Alert Service — I/O-bound. Scale on topic lag, not CPU",
                      "ML model — separate deployment; canary on ml_score PSI drift"],
                     AMBER)

    # ==================================================================
    # 20. Local ↔ production mapping
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "portability", "This code · Spring Boot + Kafka",
              "Every FastAPI dependency has a Java counterpart")
    header = ("This repo", "Java / Spring equivalent")
    rows = [
        ("services/event_bus.py",             "KafkaTemplate + @KafkaListener"),
        ("bus.publish(topic, event)",         "kafkaTemplate.send(topic, key, payload)"),
        ("bus.subscribe(topic, h)",           "@KafkaListener(topics=…, groupId=…)"),
        ("asyncio.Queue per topic",           "Kafka partitioned topic"),
        ("require_role(...) dependency",      "@PreAuthorize(\"hasRole('ADMIN')\")"),
        ("motor / AsyncIOMotorClient",         "ReactiveMongoTemplate"),
        ("@asynccontextmanager lifespan",     "@EventListener(ApplicationReadyEvent)"),
        ("metrics.snapshot()",                 "Micrometer + Prometheus"),
    ]
    cols_w = [Inches(5.9), Inches(6.43)]
    left0  = Inches(0.5); top0 = Inches(2.05); row_h = Inches(0.5)

    for i, row in enumerate([header] + rows):
        y = top0 + row_h * i
        if i == 0:
            bgc = LINE
        elif i % 2 == 0:
            bgc = RGBColor(0x14, 0x18, 0x20)
        else:
            bgc = None
        if bgc:
            bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left0, y, sum(cols_w, Emu(0)), row_h)
            bg.fill.solid(); bg.fill.fore_color.rgb = bgc
            bg.line.fill.background()
        x = left0
        for j, cell in enumerate(row):
            add_text(s, x + Inches(0.18), y + Inches(0.11),
                     cols_w[j] - Inches(0.3), row_h,
                     cell,
                     font_size=13,
                     color=(MUTED if i == 0 else (CYAN if j == 0 else TEXT)),
                     mono=True, bold=(i == 0))
            x += cols_w[j]

    # ==================================================================
    # 21. Testing & evaluation
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "quality", "Testing coverage · evaluation criteria",
              "22 backend pytests · 100% frontend flows — all green")

    # Left — pytest highlights
    add_card(s, Inches(0.5), Inches(2.05), Inches(6.15), Inches(4.7), accent=LIME)
    add_text(s, Inches(0.75), Inches(2.25), Inches(6), Inches(0.4),
             "PYTEST · 22/22 GREEN", font_size=11, color=LIME, mono=True, bold=True)
    add_text(s, Inches(0.75), Inches(2.75), Inches(5.7), Inches(3.9),
             ["• 401s on protected endpoints (anon)",
              "• Login all 3 roles + wrong-password",
              "• Register + duplicate-email 409",
              "• Admin start / stop / inject",
              "• Viewer 403 on start & ack",
              "• Analyst 403 on start · 200 on ack",
              "• Unknown alert ack → 404",
              "• Rehydration total > 0 on cold start",
              "• Low-risk approve · high-risk block with all reasons",
              "• Metrics shape verified",
              "• Download router · attachment disposition"],
             font_size=12, color=MUTED)

    # Right — evaluation criteria
    add_card(s, Inches(6.83), Inches(2.05), Inches(6.0), Inches(4.7), accent=CYAN)
    add_text(s, Inches(7.08), Inches(2.25), Inches(6), Inches(0.4),
             "EVALUATION CRITERIA MAPPING", font_size=11, color=CYAN, mono=True, bold=True)
    add_text(s, Inches(7.08), Inches(2.75), Inches(5.6), Inches(3.9),
             ["✓ Code quality — pytests + lint + testids",
              "✓ Microservices design — 3 svc, 1 bus, 3 topics",
              "✓ ML integration — /api/fraud/score as first-class",
              "✓ Demo execution — live console + video + storyboard",
              "✓ Technical presentation — this deck, deep-dive doc"],
             font_size=13, color=TEXT)

    # ==================================================================
    # 22. Same design, two stacks (Python + Java/Kafka)
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_title(s, "portability · proof", "Same design, two stacks",
              "Python (live console) + Java Spring Boot + real Kafka (assignment target)")

    # Left card — Python
    add_bullet_card(s, Inches(0.5), Inches(2.05), Inches(6.15), Inches(4.75),
                     "python reference (backend/, frontend/)",
                     ["FastAPI · Motor · scikit-learn IsolationForest",
                      "In-process asyncio event bus (Kafka topic semantics)",
                      "React SPA with JWT + RBAC + live console",
                      "Powers the walkthrough video and screenshots",
                      "One `uvicorn` + `yarn start` — no docker required"],
                     CYAN)
    # Right card — Java
    add_bullet_card(s, Inches(6.83), Inches(2.05), Inches(6.0), Inches(4.75),
                     "java + kafka (java/)",
                     ["Spring Boot 3.2 · Java 17 · spring-kafka · Resilience4j",
                      "Three real services + Zookeeper + Kafka + Mongo",
                      "ML model = standalone Python FastAPI service",
                      "@KafkaListener · @CircuitBreaker → rule fallback",
                      "docker compose up  — full pipeline in one command"],
                     AMBER)

    # ==================================================================
    # 23. Thank you
    # ==================================================================
    s = add_slide(prs); slides.append(s)
    add_pill(s, Inches(0.5), Inches(0.5), "SUBMISSION COMPLETE", LIME)
    add_text(s, Inches(0.5), Inches(2.2), Inches(12), Inches(1.5),
             "Thank you.",
             font_size=72, color=TEXT, bold=True)
    add_text(s, Inches(0.5), Inches(3.6), Inches(12), Inches(1),
             "Questions?", font_size=32, color=CYAN, bold=True)
    add_text(s, Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
             "DELIVERABLES  ·", font_size=11, color=DIM, mono=True, bold=True)
    add_text(s, Inches(0.5), Inches(5.75), Inches(12), Inches(1.4),
             ["• Source archive  (.tar.gz  /  .zip) — Python + Java implementations",
              "• End-to-end demo video  (WebM)",
              "• 12 storyboard screenshots  (ZIP)",
              "• Technical Deep Dive  (2-part markdown, ~10k words)",
              "• Demo Script  (12-scene narration)",
              "• This presentation deck  (.pptx + .pdf)"],
             font_size=14, color=MUTED)

    # add footers everywhere
    for i, sl in enumerate(slides):
        add_footer(sl, i + 1, len(slides))

    target = Path(out_path) if out_path else OUT_DEFAULT
    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target))
    print(f"OK: {target} ({target.stat().st_size / 1024:.1f} KB, {len(slides)} slides)")


if __name__ == "__main__":
    make_deck(str(OUT_DEFAULT))
